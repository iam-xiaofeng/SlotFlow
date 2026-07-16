"""Contracts and real local-format fixtures for the MarkItDown tool."""

from __future__ import annotations

import io
import json
import zipfile
from dataclasses import replace
from pathlib import Path
from types import SimpleNamespace

import fitz
import pytest
from docx import Document
from langchain_core.language_models.fake_chat_models import FakeListChatModel
from openpyxl import Workbook
from PIL import Image
from pptx import Presentation

from app.chat.runtime.config import load_markitdown_config_from_env
from app.harness.sandbox import SlotFlowSandboxConfig
from app.harness.tools import markitdown as module
from app.harness.tools.markitdown import (
    LangChainVisionClient,
    MarkItDownConversionError,
    SlotFlowMarkItDownConfig,
    build_markitdown_tools,
    convert_file_to_markdown,
)


class FakeVisionClient:
    def __init__(self, text: str = "SCANNED OCR TEXT") -> None:
        self.calls: list[dict] = []
        self.text = text
        self.chat = SimpleNamespace(completions=self)

    def create(self, **kwargs):
        self.calls.append(kwargs)
        return SimpleNamespace(
            choices=[SimpleNamespace(message=SimpleNamespace(content=self.text))],
        )


class FailingVisionClient:
    def __init__(self) -> None:
        self.chat = SimpleNamespace(completions=self)

    def create(self, **kwargs):
        del kwargs
        raise RuntimeError("synthetic Vision failure")


def _config(**updates) -> SlotFlowMarkItDownConfig:
    return replace(SlotFlowMarkItDownConfig(), **updates)


def _write_png(path: Path) -> None:
    image = Image.new("RGB", (200, 80), "white")
    image.save(path, format="PNG")


def _write_scanned_pdf(path: Path) -> None:
    buffer = io.BytesIO()
    Image.new("RGB", (300, 120), "white").save(buffer, format="PNG")
    document = fitz.open()
    page = document.new_page(width=300, height=120)
    page.insert_image(page.rect, stream=buffer.getvalue())
    document.save(path)
    document.close()


def test_disabled_config_exposes_no_conversion_tool(tmp_path: Path) -> None:
    assert build_markitdown_tools(
        _config(enabled=False),
        sandbox_config=SlotFlowSandboxConfig(workspace_root=tmp_path),
        model=None,
        thread_id="thread_docs",
    ) == []


def test_convert_local_text_and_tool_inline_projection(tmp_path: Path) -> None:
    source = tmp_path / "note.txt"
    source.write_text("# Hello\n\nLocal content")
    tools = build_markitdown_tools(
        SlotFlowMarkItDownConfig(),
        sandbox_config=SlotFlowSandboxConfig(workspace_root=tmp_path),
        model=None,
        thread_id="thread_docs",
    )

    assert [tool.name for tool in tools] == ["convert_file_to_markdown"]
    payload = json.loads(tools[0].invoke({"path": "note.txt", "use_vision": False}))
    assert payload["source_path"] == "note.txt"
    assert "Local content" in payload["markdown"]
    assert payload["vision_used"] is False


def test_tool_writes_large_result_only_inside_thread_artifacts(tmp_path: Path) -> None:
    (tmp_path / "report.html").write_text("<h1>Quarterly Report</h1><p>Revenue</p>")
    tool = build_markitdown_tools(
        SlotFlowMarkItDownConfig(),
        sandbox_config=SlotFlowSandboxConfig(workspace_root=tmp_path),
        model=None,
        thread_id="thread_docs",
    )[0]

    payload = json.loads(
        tool.invoke(
            {
                "path": "report.html",
                "output_path": "converted/report.md",
                "use_vision": False,
            }
        )
    )
    assert payload["output_path"] == "artifacts/thread_docs/converted/report.md"
    target = tmp_path / payload["output_path"]
    assert target.is_file()
    assert "Quarterly Report" in target.read_text()
    assert "markdown" not in payload


def test_tool_rejects_workspace_escape_and_input_limit(tmp_path: Path) -> None:
    outside = tmp_path.parent / "outside.txt"
    outside.write_text("secret")
    source = tmp_path / "large.txt"
    source.write_text("x" * 20)
    tool = build_markitdown_tools(
        _config(max_input_bytes=10),
        sandbox_config=SlotFlowSandboxConfig(workspace_root=tmp_path),
        model=None,
        thread_id="thread_docs",
    )[0]

    with pytest.raises(ValueError, match="must not contain"):
        tool.invoke({"path": "../outside.txt"})
    with pytest.raises(MarkItDownConversionError, match="max_input_bytes"):
        tool.invoke({"path": "large.txt"})


def test_archive_bomb_limits_are_checked_before_conversion(tmp_path: Path) -> None:
    archive = tmp_path / "bundle.zip"
    with zipfile.ZipFile(archive, "w") as zipped:
        zipped.writestr("a.txt", "a" * 20)
        zipped.writestr("b.txt", "b" * 20)

    with pytest.raises(MarkItDownConversionError, match="too many entries"):
        convert_file_to_markdown(
            archive,
            config=_config(max_archive_entries=1),
            use_vision=False,
        )
    with pytest.raises(MarkItDownConversionError, match="uncompressed size"):
        convert_file_to_markdown(
            archive,
            config=_config(max_archive_uncompressed_bytes=10),
            use_vision=False,
        )


def test_generic_zip_disables_vision_and_rejects_nested_zip(tmp_path: Path) -> None:
    safe = tmp_path / "safe.zip"
    with zipfile.ZipFile(safe, "w") as archive:
        archive.writestr("inside.txt", "Archive text")
    result = convert_file_to_markdown(
        safe,
        config=SlotFlowMarkItDownConfig(),
        llm_client=FakeVisionClient(),
        llm_model="vision-test",
    )
    assert result.vision_used is False
    assert any("generic ZIP" in warning for warning in result.warnings)

    nested_bytes = io.BytesIO()
    with zipfile.ZipFile(nested_bytes, "w") as nested:
        nested.writestr("payload.txt", "payload")
    outer = tmp_path / "outer.zip"
    with zipfile.ZipFile(outer, "w") as archive:
        archive.writestr("nested.zip", nested_bytes.getvalue())
    with pytest.raises(MarkItDownConversionError, match="nested ZIP"):
        convert_file_to_markdown(
            outer,
            config=SlotFlowMarkItDownConfig(),
            use_vision=False,
        )


def test_image_without_vision_client_returns_explicit_warning(tmp_path: Path) -> None:
    image = tmp_path / "scan.png"
    _write_png(image)

    result = convert_file_to_markdown(
        image,
        config=SlotFlowMarkItDownConfig(),
        use_vision=True,
    )

    assert result.vision_used is False
    assert any("no compatible Vision client" in warning for warning in result.warnings)


def test_pure_image_uses_openai_compatible_vision_client(tmp_path: Path) -> None:
    image = tmp_path / "scan.png"
    _write_png(image)
    client = FakeVisionClient("Invoice total: $42")

    result = convert_file_to_markdown(
        image,
        config=SlotFlowMarkItDownConfig(),
        use_vision=True,
        llm_client=client,
        llm_model="vision-test",
    )

    assert result.vision_used is True
    assert "Invoice total: $42" in result.markdown
    assert client.calls[0]["model"] == "vision-test"
    image_url = client.calls[0]["messages"][0]["content"][1]["image_url"]["url"]
    assert image_url.startswith("data:image/png;base64,")


def test_selected_vision_capable_run_model_is_reused_automatically(
    monkeypatch,
    tmp_path: Path,
) -> None:
    image = tmp_path / "scan.png"
    _write_png(image)
    model = FakeListChatModel(responses=["automatic model OCR"])
    monkeypatch.setattr(module, "_chat_model_id", lambda model: "vision-model")
    monkeypatch.setattr(module, "_model_supports_vision", lambda model_id, model: True)

    tool = build_markitdown_tools(
        SlotFlowMarkItDownConfig(),
        sandbox_config=SlotFlowSandboxConfig(workspace_root=tmp_path),
        model=model,
        thread_id="thread_docs",
    )[0]
    payload = json.loads(tool.invoke({"path": "scan.png"}))

    assert payload["vision_used"] is True
    assert "automatic model OCR" in payload["markdown"]


def test_scanned_pdf_uses_official_ocr_plugin(tmp_path: Path) -> None:
    scanned = tmp_path / "scanned.pdf"
    _write_scanned_pdf(scanned)
    client = FakeVisionClient()

    result = convert_file_to_markdown(
        scanned,
        config=SlotFlowMarkItDownConfig(),
        use_vision=True,
        llm_client=client,
        llm_model="vision-test",
    )

    assert result.vision_used is True
    assert "SCANNED OCR TEXT" in result.markdown
    assert client.calls, "markitdown-ocr must call the configured Vision client"


def test_scanned_pdf_reports_swallowed_plugin_vision_failure(tmp_path: Path) -> None:
    scanned = tmp_path / "scanned.pdf"
    _write_scanned_pdf(scanned)

    result = convert_file_to_markdown(
        scanned,
        config=SlotFlowMarkItDownConfig(),
        llm_client=FailingVisionClient(),
        llm_model="vision-test",
    )

    assert result.vision_used is True
    assert any("Vision OCR failed" in warning for warning in result.warnings)


def test_vision_page_and_embedded_image_limits(tmp_path: Path) -> None:
    pdf = fitz.open()
    pdf.new_page()
    pdf.new_page()
    pdf_path = tmp_path / "many-pages.pdf"
    pdf.save(pdf_path)
    pdf.close()

    with pytest.raises(MarkItDownConversionError, match="too many pages"):
        convert_file_to_markdown(
            pdf_path,
            config=_config(vision_max_pages=1),
            llm_client=FakeVisionClient(),
            llm_model="vision-test",
        )

    document = tmp_path / "many-images.docx"
    with zipfile.ZipFile(document, "w") as archive:
        archive.writestr("word/media/image1.png", b"a")
        archive.writestr("word/media/image2.png", b"b")
    with pytest.raises(MarkItDownConversionError, match="too many embedded images"):
        convert_file_to_markdown(
            document,
            config=_config(vision_max_images=1),
            llm_client=FakeVisionClient(),
            llm_model="vision-test",
        )


def test_real_format_matrix_uses_upstream_converters(tmp_path: Path) -> None:
    samples: list[tuple[Path, str]] = []

    html = tmp_path / "sample.html"
    html.write_text("<h1>HTML Heading</h1><p>HTML body</p>")
    samples.append((html, "HTML Heading"))

    csv = tmp_path / "sample.csv"
    csv.write_text("name,value\nalpha,7\n")
    samples.append((csv, "alpha"))

    docx = tmp_path / "sample.docx"
    word = Document()
    word.add_heading("Word Heading")
    word.add_paragraph("Word body")
    word.save(docx)
    samples.append((docx, "Word Heading"))

    xlsx = tmp_path / "sample.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.append(["Metric", "Value"])
    sheet.append(["Revenue", 42])
    workbook.save(xlsx)
    samples.append((xlsx, "Revenue"))

    pptx = tmp_path / "sample.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Slide Heading"
    slide.placeholders[1].text = "Slide body"
    presentation.save(pptx)
    samples.append((pptx, "Slide Heading"))

    pdf_path = tmp_path / "sample.pdf"
    pdf = fitz.open()
    page = pdf.new_page()
    page.insert_text((72, 72), "PDF body text")
    pdf.save(pdf_path)
    pdf.close()
    samples.append((pdf_path, "PDF body text"))

    archive = tmp_path / "sample.zip"
    with zipfile.ZipFile(archive, "w") as zipped:
        zipped.writestr("inside.txt", "Archive body")
    samples.append((archive, "Archive body"))

    for source, expected in samples:
        result = convert_file_to_markdown(
            source,
            config=SlotFlowMarkItDownConfig(),
            use_vision=False,
        )
        assert expected in result.markdown, source.name


def test_langchain_model_facade_returns_openai_shape() -> None:
    client = LangChainVisionClient(FakeListChatModel(responses=["recognized text"]))
    response = client.chat.completions.create(
        model="ignored",
        messages=[{"role": "user", "content": "read image"}],
    )
    assert response.choices[0].message.content == "recognized text"


def test_runtime_env_loads_dedicated_vision_settings(monkeypatch) -> None:
    monkeypatch.setenv("SLOTFLOW_MARKITDOWN_ENABLED", "false")
    monkeypatch.setenv("SLOTFLOW_MARKITDOWN_MAX_INPUT_BYTES", "1000")
    monkeypatch.setenv("SLOTFLOW_MARKITDOWN_VISION_MODEL", "vision-model")
    monkeypatch.setenv("SLOTFLOW_MARKITDOWN_VISION_BASE_URL", "https://vision.example/v1")
    monkeypatch.setenv("SLOTFLOW_MARKITDOWN_VISION_API_KEY", "secret")
    monkeypatch.setenv("SLOTFLOW_MARKITDOWN_VISION_MAX_PAGES", "3")

    config = load_markitdown_config_from_env()
    assert config.enabled is False
    assert config.max_input_bytes == 1000
    assert config.vision_model == "vision-model"
    assert config.vision_base_url == "https://vision.example/v1"
    assert config.vision_api_key == "secret"
    assert config.vision_max_pages == 3
